"""Generate the Section 4 slide deck (7-day, non-technical audience).

Diagrams are drawn with matplotlib and embedded as PNGs. Deck is
written to `Section_04_LLM_AI_Integration.pptx` in the same folder.

Run:
    python build_slides.py
"""

from pathlib import Path

import matplotlib.pyplot as plt
from matplotlib.patches import FancyArrow, FancyBboxPatch, Rectangle
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

HERE = Path(__file__).parent
IMG_DIR = HERE / "images"
IMG_DIR.mkdir(exist_ok=True)

# --------------------------------------------------------------------- theme
COL_BG = RGBColor(0x0E, 0x14, 0x22)
COL_ACC = RGBColor(0x36, 0xC9, 0xFF)
COL_TXT = RGBColor(0xEA, 0xEE, 0xF5)
COL_MUT = RGBColor(0x9A, 0xA3, 0xB2)

WIDE = Inches(13.333)
TALL = Inches(7.5)


def _mpl_style():
    plt.rcParams.update({
        "figure.facecolor": "#0E1422",
        "axes.facecolor":   "#0E1422",
        "axes.edgecolor":   "#9AA3B2",
        "axes.labelcolor":  "#EAEEF5",
        "xtick.color":      "#EAEEF5",
        "ytick.color":      "#EAEEF5",
        "text.color":       "#EAEEF5",
        "font.size":        14,
    })


def _box(ax, x, y, w, h, text, fc="#1B2540", ec="#36C9FF", fontsize=13):
    ax.add_patch(FancyBboxPatch((x, y), w, h,
                                boxstyle="round,pad=0.02,rounding_size=0.15",
                                fc=fc, ec=ec, lw=2))
    ax.text(x + w / 2, y + h / 2, text, ha="center", va="center",
            fontsize=fontsize, color="#EAEEF5")


def _arrow(ax, x1, y1, x2, y2, color="#36C9FF"):
    ax.add_patch(FancyArrow(x1, y1, x2 - x1, y2 - y1,
                            width=0.02, head_width=0.18, head_length=0.15,
                            length_includes_head=True, fc=color, ec=color))


# --------------------------------------------------------------------- diagrams
def img_llm_intuition(path):
    _mpl_style()
    fig, ax = plt.subplots(figsize=(11, 5))
    ax.set_xlim(0, 12); ax.set_ylim(0, 5); ax.axis("off")

    _box(ax, 0.3, 2.0, 2.4, 1.0, "You typed:\n'The cat sat on the'")
    _box(ax, 3.2, 2.0, 2.6, 1.0, "LLM looks at every\nword so far",
         fc="#231A3D", ec="#F0B000")
    _box(ax, 6.4, 2.0, 2.6, 1.0, "Guesses the most\nlikely next word",
         fc="#1B3020", ec="#3ADB90")
    _box(ax, 9.7, 2.0, 1.9, 1.0, "'mat'")

    for x1, x2 in [(2.7, 3.2), (5.8, 6.4), (9.0, 9.7)]:
        _arrow(ax, x1, 2.5, x2, 2.5)

    ax.text(6.0, 0.6, "That's it. Very fancy autocomplete, repeated for every word.",
            ha="center", fontsize=13, style="italic", color="#9AA3B2")

    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches="tight", facecolor="#0E1422")
    plt.close(fig)


def img_tokens(path):
    _mpl_style()
    fig, ax = plt.subplots(figsize=(11, 4.5))
    ax.set_xlim(0, 12); ax.set_ylim(0, 4); ax.axis("off")

    _box(ax, 0.5, 2.4, 2.0, 0.9, '"cat"',        fc="#1B3020", ec="#3ADB90", fontsize=15)
    ax.text(1.5, 2.0, "1 token", ha="center", fontsize=12, color="#9AA3B2")

    _box(ax, 3.0, 2.4, 1.4, 0.9, '"un"',         fc="#1B2540", ec="#36C9FF", fontsize=15)
    _box(ax, 4.5, 2.4, 1.7, 0.9, '"believ"',     fc="#1B2540", ec="#36C9FF", fontsize=15)
    _box(ax, 6.3, 2.4, 1.4, 0.9, '"able"',       fc="#1B2540", ec="#36C9FF", fontsize=15)
    ax.text(5.4, 2.0, "3 tokens ('unbelievable')", ha="center", fontsize=12, color="#9AA3B2")

    _box(ax, 8.2, 2.4, 1.1, 0.9, '"party"',      fc="#231A3D", ec="#F0B000", fontsize=14)
    _box(ax, 9.4, 2.4, 1.1, 0.9, '"emoji"',      fc="#231A3D", ec="#F0B000", fontsize=14)
    _box(ax, 10.6, 2.4, 1.1, 0.9, '"!"',         fc="#231A3D", ec="#F0B000", fontsize=14)
    ax.text(9.9, 2.0, "one emoji = often 2-3 tokens", ha="center", fontsize=12, color="#9AA3B2")

    ax.text(6.0, 0.7, "Rule of thumb: 1 token = about 3/4 of a word. You pay per token.",
            ha="center", fontsize=14, style="italic", color="#EAEEF5")

    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches="tight", facecolor="#0E1422")
    plt.close(fig)


def img_together_landing(path):
    _mpl_style()
    fig, ax = plt.subplots(figsize=(11, 5))
    ax.set_xlim(0, 12); ax.set_ylim(0, 5); ax.axis("off")

    _box(ax, 4.8, 2.1, 2.4, 1.0, "Your Python code",
         fc="#1B2540", ec="#36C9FF", fontsize=15)
    _box(ax, 4.8, 3.7, 2.4, 0.9, "Together AI API",
         fc="#231A3D", ec="#F0B000", fontsize=15)

    _box(ax, 0.6, 0.4, 2.2, 0.9, "LLaMA 3.1 8B",  fc="#1B3020", ec="#3ADB90", fontsize=12)
    _box(ax, 3.0, 0.4, 2.2, 0.9, "LLaMA 3.1 70B", fc="#1B3020", ec="#3ADB90", fontsize=12)
    _box(ax, 5.4, 0.4, 2.2, 0.9, "Mistral 7B",    fc="#1B3020", ec="#3ADB90", fontsize=12)
    _box(ax, 7.8, 0.4, 2.2, 0.9, "Qwen 2.5",      fc="#1B3020", ec="#3ADB90", fontsize=12)
    _box(ax, 10.2, 0.4, 1.6, 0.9, "DeepSeek",     fc="#1B3020", ec="#3ADB90", fontsize=12)

    _arrow(ax, 6.0, 2.1, 6.0, 3.7)
    _arrow(ax, 6.0, 3.7, 6.0, 2.1)
    for x in [1.7, 4.1, 6.5, 8.9, 11.0]:
        _arrow(ax, 6.0, 3.7, x, 1.3, color="#F0B000")

    ax.text(6.0, 4.9, "One API. Dozens of open-source models. No GPU needed.",
            ha="center", fontsize=14, style="italic", color="#EAEEF5")

    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches="tight", facecolor="#0E1422")
    plt.close(fig)


def img_prompt_layers(path):
    _mpl_style()
    fig, ax = plt.subplots(figsize=(10, 5.5))
    ax.set_xlim(0, 10); ax.set_ylim(0, 5.5); ax.axis("off")

    _box(ax, 1.0, 4.3, 8.0, 0.8, "ROLE — 'You are a legal analyst'",
         fc="#1B2540", ec="#36C9FF", fontsize=14)
    _box(ax, 1.0, 3.2, 8.0, 0.8, "TASK — 'Summarize this contract'",
         fc="#231A3D", ec="#F0B000", fontsize=14)
    _box(ax, 1.0, 2.1, 8.0, 0.8, "FORMAT — '3 bullets, max 15 words each'",
         fc="#1B3020", ec="#3ADB90", fontsize=14)
    _box(ax, 1.0, 1.0, 8.0, 0.8, "CONSTRAINTS — 'no legal jargon'",
         fc="#3B2222", ec="#FF7A7A", fontsize=14)

    ax.text(5.0, 0.2,
            "The gap between a weak and strong prompt is often 10x better output.",
            ha="center", fontsize=13, style="italic", color="#9AA3B2")

    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches="tight", facecolor="#0E1422")
    plt.close(fig)


def img_react(path):
    _mpl_style()
    fig, ax = plt.subplots(figsize=(11, 5.5))
    ax.set_xlim(0, 12); ax.set_ylim(0, 5.5); ax.axis("off")

    _box(ax, 0.5, 3.9, 11.0, 0.9, "Question: What's the population of the capital of France?",
         fc="#1B2540", ec="#36C9FF", fontsize=14)
    _box(ax, 0.5, 2.9, 5.2, 0.8, "Thought: I need the capital of France.",
         fc="#231A3D", ec="#F0B000", fontsize=13)
    _box(ax, 6.0, 2.9, 5.5, 0.8, "Action: search('capital of France')",
         fc="#1B3020", ec="#3ADB90", fontsize=13)
    _box(ax, 0.5, 1.9, 5.2, 0.8, "Observation: Paris",
         fc="#1B2540", ec="#36C9FF", fontsize=13)
    _box(ax, 6.0, 1.9, 5.5, 0.8, "Thought: Now I need Paris's population.",
         fc="#231A3D", ec="#F0B000", fontsize=13)
    _box(ax, 0.5, 0.9, 5.2, 0.8, "Action: search('population of Paris')",
         fc="#1B3020", ec="#3ADB90", fontsize=13)
    _box(ax, 6.0, 0.9, 5.5, 0.8, "Final Answer: about 2.1 million.",
         fc="#3B2222", ec="#FF7A7A", fontsize=14)

    ax.text(6.0, 0.2,
            "Every AI agent framework (LangGraph, CrewAI, AutoGen) is built on this loop.",
            ha="center", fontsize=13, style="italic", color="#9AA3B2")

    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches="tight", facecolor="#0E1422")
    plt.close(fig)


def img_tool_call(path):
    _mpl_style()
    fig, ax = plt.subplots(figsize=(12, 5))
    ax.set_xlim(0, 13); ax.set_ylim(0, 5); ax.axis("off")

    _box(ax, 0.2, 2.0, 2.0, 1.0, "User:\n'weather in Paris?'")
    _box(ax, 2.4, 2.0, 2.0, 1.0, "AI decides to\ncall a tool",
         fc="#231A3D", ec="#F0B000")
    _box(ax, 4.6, 2.0, 2.4, 1.0, "AI writes:\nget_weather(\n  city='Paris')",
         fc="#1B3020", ec="#3ADB90")
    _box(ax, 7.3, 2.0, 2.0, 1.0, "YOUR code\nruns the fn")
    _box(ax, 9.6, 2.0, 1.6, 1.0, "Result:\n21C, sunny",
         fc="#1B3020", ec="#3ADB90")
    _box(ax, 11.4, 2.0, 1.5, 1.0, "AI writes\nfinal answer",
         fc="#231A3D", ec="#F0B000")

    for x1, x2 in [(2.2, 2.4), (4.4, 4.6), (7.0, 7.3), (9.3, 9.6), (11.2, 11.4)]:
        _arrow(ax, x1, 2.5, x2, 2.5)

    ax.text(6.5, 0.6,
            "AI never runs code. It writes a request. YOU decide to honor it.",
            ha="center", fontsize=13, style="italic", color="#9AA3B2")

    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches="tight", facecolor="#0E1422")
    plt.close(fig)


def img_stream_vs_block(path):
    _mpl_style()
    fig, ax = plt.subplots(figsize=(10, 5))
    ax.set_xlim(0, 10); ax.set_ylim(0, 5); ax.axis("off")

    ax.text(1.0, 4.4, "Blocking (no streaming)", fontsize=15, color="#EAEEF5")
    ax.add_patch(Rectangle((1.0, 3.7), 6.0, 0.4, fc="#3B2222", ec="#FF7A7A"))
    ax.text(4.0, 3.9, "wait...wait...wait...wait...", ha="center", fontsize=12)
    ax.add_patch(Rectangle((7.0, 3.7), 1.8, 0.4, fc="#3ADB90", ec="#3ADB90"))
    ax.text(7.9, 3.9, "full answer", ha="center", fontsize=12, color="#0E1422")

    ax.text(1.0, 2.6, "Streaming (stream=True)", fontsize=15, color="#EAEEF5")
    for x in [1.2, 2.0, 2.8, 3.6, 4.4, 5.2, 6.0, 6.8, 7.6, 8.4]:
        ax.add_patch(Rectangle((x, 1.9), 0.7, 0.4, fc="#3ADB90", ec="#3ADB90"))
    ax.text(4.8, 2.35, "token...token...token...token...", ha="center", fontsize=12, color="#0E1422")

    ax.text(5.0, 0.8,
            "Same total time. First word appears in ~300 ms instead of 4 seconds.",
            ha="center", fontsize=13, style="italic", color="#9AA3B2")

    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches="tight", facecolor="#0E1422")
    plt.close(fig)


def img_async(path):
    _mpl_style()
    fig, ax = plt.subplots(figsize=(10, 5))
    ax.set_xlim(0, 10); ax.set_ylim(0, 5); ax.axis("off")

    ax.text(1.0, 4.5, "Sequential (regular Python)", fontsize=15, color="#EAEEF5")
    for i, x in enumerate([1.0, 2.5, 4.0, 5.5]):
        ax.add_patch(Rectangle((x, 3.8), 1.3, 0.4, fc="#3B2222", ec="#FF7A7A"))
        ax.text(x + 0.65, 4.0, f"call {i+1}", ha="center", fontsize=11)
    ax.text(1.0, 3.4, "Total time = 4 × per-call time", fontsize=12, color="#9AA3B2")

    ax.text(1.0, 2.6, "Parallel (asyncio.gather)", fontsize=15, color="#EAEEF5")
    for i in range(4):
        ax.add_patch(Rectangle((1.0, 1.4 + i * 0.35), 1.3, 0.3,
                               fc="#3ADB90", ec="#3ADB90"))
        ax.text(1.65, 1.55 + i * 0.35, f"call {i+1}", ha="center", fontsize=10, color="#0E1422")
    ax.text(1.0, 1.0, "Total time ≈ 1 × per-call time (all at once)",
            fontsize=12, color="#9AA3B2")

    ax.text(5.0, 0.3,
            "AI calls are 99% waiting. Async lets you wait on many at once.",
            ha="center", fontsize=13, style="italic", color="#EAEEF5")

    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches="tight", facecolor="#0E1422")
    plt.close(fig)


def img_pricing(path):
    _mpl_style()
    fig, ax = plt.subplots(figsize=(10, 5))
    models = ["Together\nLLaMA 8B", "Together\nLLaMA 70B",
              "GPT-4o-mini", "GPT-4o",
              "Claude 3.5\nHaiku", "Claude 3.5\nSonnet"]
    input_cost  = [0.18, 0.88, 0.15, 2.50, 0.80, 3.00]
    output_cost = [0.18, 0.88, 0.60, 10.00, 4.00, 15.00]

    x = list(range(len(models)))
    ax.bar([i - 0.2 for i in x], input_cost, width=0.4, color="#36C9FF", label="input $/M")
    ax.bar([i + 0.2 for i in x], output_cost, width=0.4, color="#F0B000", label="output $/M")
    ax.set_xticks(x); ax.set_xticklabels(models, fontsize=11)
    ax.set_ylabel("USD per 1,000,000 tokens")
    ax.set_title("Rough LLM pricing (early 2026)", color="#EAEEF5", pad=12)
    ax.legend(facecolor="#1B2540", edgecolor="#9AA3B2", labelcolor="#EAEEF5")
    ax.grid(axis="y", color="#2A3555", alpha=0.4)

    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches="tight", facecolor="#0E1422")
    plt.close(fig)


def img_router(path):
    _mpl_style()
    fig, ax = plt.subplots(figsize=(11, 5.5))
    ax.set_xlim(0, 12); ax.set_ylim(0, 6); ax.axis("off")

    _box(ax, 0.3, 2.6, 1.8, 1.0, "User prompt")
    _box(ax, 2.4, 2.6, 2.2, 1.0, "route(prompt)",
         fc="#231A3D", ec="#F0B000", fontsize=15)

    _box(ax, 6.0, 4.4, 3.4, 0.9,
         "Short / general → Together AI ($)",
         fc="#1B3020", ec="#3ADB90", fontsize=13)
    _box(ax, 6.0, 2.8, 3.4, 0.9,
         "Code or long doc → Claude 3.5 Sonnet ($$$)",
         fc="#1B2540", ec="#36C9FF", fontsize=13)
    _box(ax, 6.0, 1.2, 3.4, 0.9,
         "Privacy mode → Together (open source)",
         fc="#3B2222", ec="#FF7A7A", fontsize=13)

    _arrow(ax, 2.1, 3.1, 2.4, 3.1)
    _arrow(ax, 4.7, 3.1, 6.0, 4.85)
    _arrow(ax, 4.7, 3.1, 6.0, 3.25)
    _arrow(ax, 4.7, 3.1, 6.0, 1.65)

    ax.text(6.0, 0.4,
            "Cheapest model that clears the task — not 'best model always'.",
            ha="center", fontsize=13, style="italic", color="#9AA3B2")

    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches="tight", facecolor="#0E1422")
    plt.close(fig)


def img_capstone(path):
    _mpl_style()
    fig, ax = plt.subplots(figsize=(12, 6))
    ax.set_xlim(0, 12); ax.set_ylim(0, 6); ax.axis("off")

    _box(ax, 0.3, 2.6, 1.6, 1.0, "Client\n(curl / UI)")
    _box(ax, 2.2, 2.6, 2.0, 1.0, "FastAPI\n+ JWT",
         fc="#231A3D", ec="#F0B000")
    _box(ax, 4.5, 2.6, 2.2, 1.0, "Router\n+ fallback",
         fc="#1B2540", ec="#36C9FF")
    _box(ax, 7.0, 4.4, 2.4, 0.9, "Together AI\nLLaMA / Mistral",
         fc="#1B3020", ec="#3ADB90")
    _box(ax, 7.0, 2.8, 2.4, 0.9, "OpenAI\nGPT-4o / mini",
         fc="#1B3020", ec="#3ADB90")
    _box(ax, 7.0, 1.2, 2.4, 0.9, "Anthropic\nClaude 3.5",
         fc="#1B3020", ec="#3ADB90")
    _box(ax, 9.9, 2.6, 2.0, 1.0, "SQLite\ncost + budget",
         fc="#3B2222", ec="#FF7A7A")

    for x1, x2 in [(1.9, 2.2), (4.2, 4.5), (6.7, 7.0), (9.4, 9.9)]:
        _arrow(ax, x1, 3.1, x2, 3.1)
    _arrow(ax, 6.7, 3.1, 7.0, 4.85)
    _arrow(ax, 6.7, 3.1, 7.0, 1.65)

    ax.text(6.0, 0.5,
            "Every response streams to the client; every call logs a row.",
            ha="center", fontsize=13, style="italic", color="#9AA3B2")

    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches="tight", facecolor="#0E1422")
    plt.close(fig)


# --------------------------------------------------------------------- slide helpers
def add_bg(slide, prs):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = COL_BG
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element); slide.shapes._spTree.insert(2, bg._element)


def add_title(slide, prs, text, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), prs.slide_width - Inches(1.1), Inches(1.0))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = COL_TXT

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.55), Inches(1.15), prs.slide_width - Inches(1.1), Inches(0.5))
        p = sub.text_frame.paragraphs[0]; p.text = subtitle
        p.font.size = Pt(16); p.font.color.rgb = COL_ACC


def add_bullets(slide, prs, bullets, top=1.9, width=12.2, height=5.0, font=19):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame; tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(font); p.font.color.rgb = COL_TXT
        p.space_after = Pt(8)


def add_note(slide, prs, note):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(6.9),
                                   prs.slide_width - Inches(1.1), Inches(0.4))
    p = box.text_frame.paragraphs[0]; p.text = note
    p.font.size = Pt(14); p.font.italic = True; p.font.color.rgb = COL_MUT


def add_image(slide, prs, path, left=1.5, top=2.0, width=10.3):
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))


# --------------------------------------------------------------------- build
def build():
    prs = Presentation()
    prs.slide_width = WIDE; prs.slide_height = TALL
    blank = prs.slide_layouts[6]

    # ---- Cover ----
    s = prs.slides.add_slide(blank); add_bg(s, prs)
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.5), Inches(2.0))
    p = box.text_frame.paragraphs[0]; p.text = "Section 4"
    p.font.size = Pt(28); p.font.color.rgb = COL_ACC
    p = box.text_frame.add_paragraph(); p.text = "LLM Engineering & AI Integration"
    p.font.size = Pt(50); p.font.bold = True; p.font.color.rgb = COL_TXT
    p = box.text_frame.add_paragraph(); p.text = "7 days · ~1 hour each · non-technical friendly"
    p.font.size = Pt(20); p.font.color.rgb = COL_MUT

    # ---- Roadmap ----
    s = prs.slides.add_slide(blank); add_bg(s, prs)
    add_title(s, prs, "7-Day Roadmap",
              "From 'what is AI' to shipping a small multi-model assistant")
    add_bullets(s, prs, [
        "Day 1 — What is an LLM? Tokens, temperature, hallucinations",
        "Day 2 — Talking to OpenAI & Claude",
        "Day 3 — Hugging Face hands-on + Together AI",
        "Day 4 — Prompt engineering: role, few-shot, CoT, ReAct",
        "Day 5 — Structured outputs & tool calling",
        "Day 6 — Cost control, streaming & async",
        "Day 7 — Capstone: your own multi-model AI Assistant",
    ], font=20)

    # ---- Day 1 ----
    img_llm_intuition(IMG_DIR / "llm.png")
    img_tokens(IMG_DIR / "tokens.png")

    s = prs.slides.add_slide(blank); add_bg(s, prs)
    add_title(s, prs, "Day 1 — What is an LLM?",
              "Answer: it's very fancy autocomplete.")
    add_image(s, prs, IMG_DIR / "llm.png", left=0.5, top=1.7, width=12.3)

    s = prs.slides.add_slide(blank); add_bg(s, prs)
    add_title(s, prs, "Day 1 — Tokens: how AI chops up text")
    add_image(s, prs, IMG_DIR / "tokens.png", left=0.6, top=1.9, width=12.1)
    add_note(s, prs, "You are billed per token by every AI provider. More tokens = more money.")

    s = prs.slides.add_slide(blank); add_bg(s, prs)
    add_title(s, prs, "Day 1 — Other key ideas")
    add_bullets(s, prs, [
        "Attention — the AI 'looks at' every earlier word to guess the next one",
        "Temperature — a dial from 'boring & factual' (0) to 'creative' (1.5)",
        "Hallucinations — LLMs invent plausible-sounding but false answers",
        "Real-world: ChatGPT, Claude, Copilot, Cursor all use these same ideas",
    ], font=22)

    # ---- Day 2 ----
    s = prs.slides.add_slide(blank); add_bg(s, prs)
    add_title(s, prs, "Day 2 — Using AI APIs",
              "OpenAI, Claude, and what Hugging Face is")
    add_bullets(s, prs, [
        "Store keys safely in a .env file — never commit them",
        "OpenAI: client.chat.completions.create(model, messages)",
        "Claude: same idea, system prompt is a separate parameter",
        "Multi-turn: send the whole history each time (LLMs are stateless)",
        "Hugging Face = 'GitHub for AI models'—Day 3 gets hands-on",
    ], font=21)

    s = prs.slides.add_slide(blank); add_bg(s, prs)
    add_title(s, prs, "Day 2 — Real-world examples")
    add_bullets(s, prs, [
        "ChatGPT — the app, uses GPT-4o under the hood",
        "GitHub Copilot — OpenAI models fine-tuned on code",
        "Notion AI — runs on Claude 3.5 Sonnet",
        "Cursor (AI code editor) — routes coding to Claude, chit-chat to a cheaper model",
        "You'll build a mini Cursor-style router on Day 7",
    ], font=22)

    # ---- Day 3 ----
    img_together_landing(IMG_DIR / "together.png")

    s = prs.slides.add_slide(blank); add_bg(s, prs)
    add_title(s, prs, "Day 3 — Hugging Face hands-on",
              "One line of code, thousands of open-source models")
    add_bullets(s, prs, [
        "pipeline('sentiment-analysis')(text) — real AI on your CPU",
        "pipeline('summarization') — turn a page into a paragraph",
        "SentenceTransformer(...) — turn text into vectors (preview of Section 5)",
        "Small models (<1 GB) run great on laptops. Big models (8-70B) don't.",
        "For big open-source models: Together AI (next slide)",
    ], font=21)

    s = prs.slides.add_slide(blank); add_bg(s, prs)
    add_title(s, prs, "Day 3 — Together AI",
              "Open-source LLMs via one API. No GPU needed.")
    add_image(s, prs, IMG_DIR / "together.png", left=0.5, top=1.9, width=12.3)

    # ---- Day 4 ----
    img_prompt_layers(IMG_DIR / "prompts.png")
    img_react(IMG_DIR / "react.png")

    s = prs.slides.add_slide(blank); add_bg(s, prs)
    add_title(s, prs, "Day 4 — Prompt Engineering",
              "Same model, better prompt = dramatically better answers")
    add_image(s, prs, IMG_DIR / "prompts.png", left=1.5, top=1.7, width=10.3)

    s = prs.slides.add_slide(blank); add_bg(s, prs)
    add_title(s, prs, "Day 4 — Two tricks that always help")
    add_bullets(s, prs, [
        "Few-shot: show 2-3 examples of the format you want",
        "Chain of Thought: add the phrase 'Let's think step by step'",
        "Wei et al. (2022) — CoT lifted math accuracy from 18% -> 57%",
        "Real-world: Cursor uses CoT before it edits your code",
    ], font=22)

    s = prs.slides.add_slide(blank); add_bg(s, prs)
    add_title(s, prs, "Day 4 — ReAct: the shape of an AI agent",
              "Thought -> Action -> Observation -> ... -> Final Answer")
    add_image(s, prs, IMG_DIR / "react.png", left=0.5, top=1.9, width=12.3)

    # ---- Day 5 ----
    img_tool_call(IMG_DIR / "tool_call.png")

    s = prs.slides.add_slide(blank); add_bg(s, prs)
    add_title(s, prs, "Day 5 — Structured outputs",
              "Turn AI text into typed JSON your code can use")
    add_bullets(s, prs, [
        "Ask for JSON in the prompt — works about 90% of the time",
        "response_format={'type':'json_object'} — guarantees valid JSON",
        "Then json.loads() safely — no crashes, no regex",
        "Real-world: Ramp auto-extracts expense data from receipts this way",
    ], font=22)

    s = prs.slides.add_slide(blank); add_bg(s, prs)
    add_title(s, prs, "Day 5 — Tool calling",
              "How AI apps take real actions in the world")
    add_image(s, prs, IMG_DIR / "tool_call.png", left=0.3, top=1.9, width=12.7)
    add_note(s, prs, "Zapier, Notion AI, Copilot's file edits — all built on this exact pattern.")

    # ---- Day 6 ----
    img_stream_vs_block(IMG_DIR / "stream.png")
    img_async(IMG_DIR / "async.png")
    img_pricing(IMG_DIR / "pricing.png")

    s = prs.slides.add_slide(blank); add_bg(s, prs)
    add_title(s, prs, "Day 6 — Cost control")
    add_image(s, prs, IMG_DIR / "pricing.png", left=1.0, top=1.7, width=11.3)
    add_note(s, prs, "Rule of thumb: output tokens cost 3-5x more than input tokens.")

    s = prs.slides.add_slide(blank); add_bg(s, prs)
    add_title(s, prs, "Day 6 — Streaming makes it feel instant")
    add_image(s, prs, IMG_DIR / "stream.png", left=1.3, top=1.9, width=10.7)
    add_note(s, prs, "Same total time — first word appears in 300 ms instead of 4 s.")

    s = prs.slides.add_slide(blank); add_bg(s, prs)
    add_title(s, prs, "Day 6 — Async: many calls at once",
              "AI calls are 99% waiting. asyncio.gather lets you wait on them together.")
    add_image(s, prs, IMG_DIR / "async.png", left=1.3, top=1.9, width=10.7)

    # ---- Day 7 ----
    img_router(IMG_DIR / "router.png")
    img_capstone(IMG_DIR / "capstone.png")

    s = prs.slides.add_slide(blank); add_bg(s, prs)
    add_title(s, prs, "Day 7 — Multi-model routing",
              "One assistant, three providers, right model per request")
    add_image(s, prs, IMG_DIR / "router.png", left=0.5, top=1.7, width=12.3)

    s = prs.slides.add_slide(blank); add_bg(s, prs)
    add_title(s, prs, "Day 7 — Capstone architecture")
    add_image(s, prs, IMG_DIR / "capstone.png", left=0.4, top=1.7, width=12.5)

    s = prs.slides.add_slide(blank); add_bg(s, prs)
    add_title(s, prs, "What you can build after Section 4")
    add_bullets(s, prs, [
        "AI chatbots with per-user budgets and provider fallback",
        "Structured extraction from receipts, resumes, medical notes",
        "AI agents that call your APIs via tool calling",
        "Private assistants that stay on open-source models for sensitive data",
        "Anything that needs to be fast, cheap AND smart",
    ], font=22)

    # ---- Closing ----
    s = prs.slides.add_slide(blank); add_bg(s, prs)
    box = s.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.5), Inches(2.5))
    p = box.text_frame.paragraphs[0]; p.text = "Section 4 complete."
    p.font.size = Pt(52); p.font.bold = True; p.font.color.rgb = COL_TXT
    p = box.text_frame.add_paragraph(); p.text = "Next — Section 5: Embeddings, Vector Search & Semantic Systems"
    p.font.size = Pt(22); p.font.color.rgb = COL_ACC

    out = HERE / "Section_04_LLM_AI_Integration.pptx"
    prs.save(out)
    print(f"OK  {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
